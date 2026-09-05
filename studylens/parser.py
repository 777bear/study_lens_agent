from __future__ import annotations

import hashlib
import os
import shutil
import subprocess
from pathlib import Path
from typing import Iterable

from .models import SlideContent


SUPPORTED_EXTENSIONS = {".pptx", ".pdf"}


class PresentationParseError(RuntimeError):
    pass


class PresentationParser:
    def __init__(self, max_file_mb: int = 80):
        self.max_file_bytes = max_file_mb * 1024 * 1024

    def parse(self, source: str | Path, workspace: str | Path) -> list[SlideContent]:
        source = Path(source).expanduser().resolve()
        workspace = Path(workspace).expanduser().resolve()
        workspace.mkdir(parents=True, exist_ok=True)
        if not source.exists() or not source.is_file():
            raise PresentationParseError("找不到上传的课件文件。")
        if source.suffix.lower() not in SUPPORTED_EXTENSIONS:
            raise PresentationParseError("当前支持 .pptx 和 .pdf 文件。旧版 .ppt 请先另存为 .pptx。")
        if source.stat().st_size > self.max_file_bytes:
            raise PresentationParseError(f"文件超过 {self.max_file_bytes // 1024 // 1024} MB 限制。")
        if source.suffix.lower() == ".pptx":
            slides = self._parse_pptx(source, workspace)
        else:
            slides = self._parse_pdf(source, workspace)
        rendered = self._render_pages(source, workspace / "rendered")
        for slide, image_path in zip(slides, rendered):
            slide.rendered_image = str(image_path)
        return slides

    @staticmethod
    def fingerprint(source: str | Path) -> str:
        digest = hashlib.sha256()
        with Path(source).open("rb") as handle:
            for chunk in iter(lambda: handle.read(1024 * 1024), b""):
                digest.update(chunk)
        return digest.hexdigest()

    def _parse_pptx(self, source: Path, workspace: Path) -> list[SlideContent]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise PresentationParseError("缺少 python-pptx，请先安装 requirements.txt。") from exc

        media_dir = workspace / "media"
        media_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation(str(source))
        slides: list[SlideContent] = []

        def iter_shapes(shapes: Iterable):
            ordered = sorted(shapes, key=lambda shape: (int(getattr(shape, "top", 0)), int(getattr(shape, "left", 0))))
            for shape in ordered:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    yield from iter_shapes(shape.shapes)
                else:
                    yield shape

        for index, slide in enumerate(presentation.slides, start=1):
            title = ""
            text_parts: list[str] = []
            tables: list[str] = []
            charts: list[str] = []
            images: list[str] = []
            if slide.shapes.title is not None:
        
                title_shape = slide.shapes.title
                title = (getattr(title_shape, "text", "") or "").strip()

            for shape_index, shape in enumerate(iter_shapes(slide.shapes), start=1):
                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                    if rows:
                        tables.append("\n".join(rows))
                    continue
                if getattr(shape, "has_chart", False):
                    chart_lines: list[str] = []
                    try:
                        for series in shape.chart.series:
                            values = list(getattr(series, "values", []) or [])
                            chart_lines.append(f"{getattr(series, 'name', '系列')}: {values}")
                    except Exception:
                        chart_lines.append("检测到图表，数值需结合页面图像理解。")
                    charts.append("\n".join(chart_lines))
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        extension = shape.image.ext or "png"
                        target = media_dir / f"slide-{index:03d}-image-{shape_index:03d}.{extension}"
                        target.write_bytes(shape.image.blob)
                        images.append(str(target))
                    except Exception:
                        pass
                    continue
                if getattr(shape, "has_text_frame", False):
                    content = (shape.text or "").strip()
                    if content and content != title:
                        text_parts.append(content)

            notes = ""
            try:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            except Exception:
                pass
            slides.append(
                SlideContent(
                    number=index,
                    title=title,
                    text="\n".join(text_parts),
                    notes=notes,
                    tables=tables,
                    charts=charts,
                    images=images,
                )
            )
        return slides

    def _parse_pdf(self, source: Path, workspace: Path) -> list[SlideContent]:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise PresentationParseError("缺少 pypdf，请先安装 requirements.txt。") from exc
        reader = PdfReader(str(source))
        slides: list[SlideContent] = []
        for index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            title = lines[0][:120] if lines else f"第{index}页"
            slides.append(SlideContent(number=index, title=title, text="\n".join(lines[1:])))
        return slides

    def _render_pages(self, source: Path, output_dir: Path) -> list[Path]:
        output_dir.mkdir(parents=True, exist_ok=True)
        pdf_path = source
        if source.suffix.lower() == ".pptx":
            office = shutil.which("soffice") or shutil.which("libreoffice")
            if not office:
                return []
            converted_dir = output_dir / "converted"
            converted_dir.mkdir(exist_ok=True)
            env = os.environ.copy()
            env["HOME"] = str(output_dir / "lo-home")
            subprocess.run(
                [office, "--headless", "--convert-to", "pdf", "--outdir", str(converted_dir), str(source)],
                check=False,
                capture_output=True,
                timeout=120,
                env=env,
            )
            candidates = sorted(converted_dir.glob("*.pdf"))
            if not candidates:
                return []
            pdf_path = candidates[0]

        try:
            import fitz

            document = fitz.open(str(pdf_path))
            paths: list[Path] = []
            for index, page in enumerate(document, start=1):
                target = output_dir / f"slide-{index:03d}.png"
                page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False).save(str(target))
                paths.append(target)
            document.close()
            return paths
        except ImportError:
            pass

        pdftoppm = shutil.which("pdftoppm")
        if not pdftoppm:
            return []
        prefix = output_dir / "slide"
        subprocess.run(
            [pdftoppm, "-png", "-r", "144", str(pdf_path), str(prefix)],
            check=False,
            capture_output=True,
            timeout=180,
        )
        pages = sorted(output_dir.glob("slide-*.png"), key=lambda path: int(path.stem.rsplit("-", 1)[-1]))
        normalized: list[Path] = []
        for index, page in enumerate(pages, start=1):
            target = output_dir / f"slide-{index:03d}.png"
            if page != target:
                page.replace(target)
            normalized.append(target)
        return normalized
